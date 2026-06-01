#!/usr/bin/env python3
"""
Genesis AI Micro-Game Engine - Upgraded Presentation Generator
Generates a comprehensive 27-slide PPTX with dark theme and brand colors.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Brand Colors ──────────────────────────────────────────────
BG_DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # #1A1A2E
ACCENT_RED = RGBColor(0xCF, 0x0A, 0x2C)   # #CF0A2C
GOLD       = RGBColor(0xFF, 0xD7, 0x00)   # #FFD700
TEAL       = RGBColor(0x4E, 0xCD, 0xC4)   # #4ECDC4
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xCC)
DARK_PANEL = RGBColor(0x22, 0x22, 0x3E)
DARK_PANEL2= RGBColor(0x16, 0x16, 0x28)
MEDIUM_GRAY= RGBColor(0x88, 0x88, 0xAA)

# ── Helpers ───────────────────────────────────────────────────
def set_slide_bg(slide, color=BG_DARK):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Add a shape with optional fill and border."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14, color=WHITE,
                bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑'):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=12, color=WHITE,
                          bold=False, alignment=PP_ALIGN.LEFT, font_name='微软雅黑', line_spacing=1.2):
    """Add a text box with multiple lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line_data, dict):
            p.text = line_data.get('text', '')
            p.font.size = Pt(line_data.get('size', font_size))
            p.font.color.rgb = line_data.get('color', color)
            p.font.bold = line_data.get('bold', bold)
            p.font.name = line_data.get('font', font_name)
            p.alignment = line_data.get('align', alignment)
            if 'space_before' in line_data:
                p.space_before = Pt(line_data['space_before'])
            if 'space_after' in line_data:
                p.space_after = Pt(line_data['space_after'])
        else:
            p.text = str(line_data)
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.font.name = font_name
            p.alignment = alignment
    return txBox

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None,
              header_color=ACCENT_RED, row_colors=(DARK_PANEL, DARK_PANEL2)):
    """Add a styled table."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ''
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.fore_color.rgb = row_colors[r % 2]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = WHITE if r == 0 else LIGHT_GRAY
                paragraph.font.name = '微软雅黑'
                paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape

def add_accent_line(slide, left, top, width, color=ACCENT_RED, height=Pt(3)):
    """Add a horizontal accent line."""
    shape = add_shape(slide, left, top, width, height, fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    return shape

def add_card(slide, left, top, width, height, title, body_lines, title_color=GOLD, icon_text=""):
    """Add a styled card with title and body text."""
    card = add_shape(slide, left, top, width, height, fill_color=DARK_PANEL, line_color=RGBColor(0x33, 0x33, 0x55), line_width=Pt(1))
    
    if icon_text:
        add_textbox(slide, left + Inches(0.15), top + Inches(0.1), Inches(0.5), Inches(0.4),
                    icon_text, font_size=18, color=TEAL, bold=True)
        title_left = left + Inches(0.6)
    else:
        title_left = left + Inches(0.15)
    
    add_textbox(slide, title_left, top + Inches(0.1), width - Inches(0.3), Inches(0.35),
                title, font_size=13, color=title_color, bold=True)
    
    lines = []
    for bl in body_lines:
        lines.append({'text': bl, 'size': 10, 'color': LIGHT_GRAY})
    
    if lines:
        add_multiline_textbox(slide, left + Inches(0.15), top + Inches(0.5), width - Inches(0.3), height - Inches(0.6),
                              lines, font_size=10, color=LIGHT_GRAY)
    return card

def add_section_header(slide, part_num, part_title, subtitle=""):
    """Add a section divider slide."""
    set_slide_bg(slide)
    # Large part number
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(1.5),
                f"PART {part_num}", font_size=48, color=RGBColor(0x33, 0x33, 0x55), bold=True, alignment=PP_ALIGN.LEFT)
    add_accent_line(slide, Inches(0.8), Inches(2.8), Inches(2), color=ACCENT_RED)
    add_textbox(slide, Inches(0.8), Inches(3.0), Inches(8), Inches(1),
                part_title, font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(4.0), Inches(8), Inches(0.8),
                    subtitle, font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

def add_slide_number(slide, num, total=27):
    """Add slide number at bottom right."""
    add_textbox(slide, Inches(8.5), Inches(7.0), Inches(1.2), Inches(0.3),
                f"{num} / {total}", font_size=9, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)

def add_footer(slide, slide_num, total=27):
    """Add footer with branding and slide number."""
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(3), Inches(0.3),
                "Genesis — AI微游戏生成器", font_size=8, color=MEDIUM_GRAY)
    add_slide_number(slide, slide_num, total)

def add_arrow_box(slide, left, top, width, height, text, color=TEAL, text_color=WHITE):
    """Add an arrow-shaped box."""
    shape = add_shape(slide, left, top, width, height, fill_color=color,
                      shape_type=MSO_SHAPE.CHEVRON)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = '微软雅黑'
    p.alignment = PP_ALIGN.CENTER
    return shape


# ══════════════════════════════════════════════════════════════
# MAIN PRESENTATION BUILDER
# ══════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

# ────────────────────────────────────────────────────────────
# SLIDE 1: COVER
# ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Top accent bar
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), fill_color=ACCENT_RED, shape_type=MSO_SHAPE.RECTANGLE)

# Genesis logo area
add_shape(slide, Inches(3.5), Inches(1.2), Inches(3), Inches(0.6), fill_color=None,
          line_color=ACCENT_RED, line_width=Pt(2), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, Inches(3.5), Inches(1.25), Inches(3), Inches(0.5),
            "◆  GENESIS  ◆", font_size=20, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)

# Main title
add_textbox(slide, Inches(1), Inches(2.2), Inches(8), Inches(1.2),
            "AI 微游戏生成器", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide, Inches(1.5), Inches(3.5), Inches(7), Inches(0.6),
            "华为游戏中心 × 鸿蒙元服务  全方案", font_size=22, color=GOLD, bold=False, alignment=PP_ALIGN.CENTER)

# Accent line
add_accent_line(slide, Inches(3.5), Inches(4.3), Inches(3), color=ACCENT_RED, height=Pt(2))

# Version
add_textbox(slide, Inches(1), Inches(4.8), Inches(8), Inches(0.5),
            "Version 2.0 MVP  |  全方位升级方案", font_size=16, color=TEAL, alignment=PP_ALIGN.CENTER)

# Key metrics bar
metrics_y = Inches(5.6)
for i, (val, label) in enumerate([("15", "游戏模板"), ("9", "智能触发器"), ("363+", "测试用例"), ("94%+", "覆盖率")]):
    x = Inches(1.5) + Inches(1.8) * i
    add_textbox(slide, x, metrics_y, Inches(1.5), Inches(0.4),
                val, font_size=28, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, metrics_y + Inches(0.4), Inches(1.5), Inches(0.3),
                label, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom tag
add_textbox(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.4),
            "让每个等待时刻都变成游戏时刻", font_size=12, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

print("✓ Slide 1: Cover")

# ────────────────────────────────────────────────────────────
# SLIDE 2: 目录
# ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 2)

add_textbox(slide, Inches(0.6), Inches(0.4), Inches(5), Inches(0.6),
            "目录  CONTENTS", font_size=28, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(1.0), Inches(1.5), color=ACCENT_RED)

toc_items = [
    ("01", "产品概览", "定位、价值主张、用户场景", ACCENT_RED),
    ("02", "技术架构", "分层架构、核心引擎、游戏模板、触发器", ACCENT_RED),
    ("03", "AI 能力", "画像系统、推荐引擎、动态难度、场景感知", ACCENT_RED),
    ("04", "鸿蒙集成", "元服务架构、服务卡片、ArkTS桥接", ACCENT_RED),
    ("05", "实现成果", "工程成果、Demo展示、CI/CD、性能指标", ACCENT_RED),
    ("06", "发布计划", "里程碑、上架流程、监控运维、商业化", ACCENT_RED),
]

for i, (num, title, desc, color) in enumerate(toc_items):
    y = Inches(1.4) + Inches(0.9) * i
    # Number circle
    add_shape(slide, Inches(0.8), y, Inches(0.6), Inches(0.6), fill_color=color, shape_type=MSO_SHAPE.OVAL)
    add_textbox(slide, Inches(0.8), y + Inches(0.08), Inches(0.6), Inches(0.45),
                num, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.6), y + Inches(0.02), Inches(3), Inches(0.35),
                title, font_size=18, color=WHITE, bold=True)
    add_textbox(slide, Inches(1.6), y + Inches(0.35), Inches(6), Inches(0.3),
                desc, font_size=11, color=LIGHT_GRAY)
    # Decorative line
    add_shape(slide, Inches(1.6), y + Inches(0.65), Inches(6.5), Pt(0.5),
              fill_color=RGBColor(0x33, 0x33, 0x55), shape_type=MSO_SHAPE.RECTANGLE)

print("✓ Slide 2: Contents")

# ────────────────────────────────────────────────────────────
# PART 1: 产品概览
# ────────────────────────────────────────────────────────────

# SLIDE 3: Section Divider - Part 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "01", "产品概览", "Product Overview")
add_slide_number(slide, 3)
print("✓ Slide 3: Part 1 Section")

# SLIDE 4: 产品定位与愿景
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 4)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "产品定位与愿景", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Vision quote
add_shape(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(1.2), fill_color=DARK_PANEL,
          line_color=GOLD, line_width=Pt(1.5))
add_textbox(slide, Inches(1.0), Inches(1.3), Inches(8), Inches(1.0),
            '"让每个等待时刻都变成游戏时刻"', font_size=24, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# Three pillars
pillars = [
    ("🎮", "微游戏引擎", "15种即开即玩的小游戏\n平均<100KB，零安装"),
    ("🤖", "AI智能驱动", "用户画像 + 场景感知\n千人千面的个性化体验"),
    ("📱", "鸿蒙元服务", "免安装服务卡片\n无缝嵌入华为游戏中心"),
]
for i, (icon, title, desc) in enumerate(pillars):
    x = Inches(0.6) + Inches(3.1) * i
    card = add_shape(slide, x, Inches(2.8), Inches(2.8), Inches(2.8),
                     fill_color=DARK_PANEL, line_color=TEAL, line_width=Pt(1))
    add_textbox(slide, x + Inches(0.3), Inches(2.95), Inches(2.2), Inches(0.5),
                icon, font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(3.5), Inches(2.4), Inches(0.4),
                title, font_size=16, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(4.1), Inches(2.4), Inches(1.2),
                desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom note
add_textbox(slide, Inches(0.6), Inches(6.1), Inches(8.8), Inches(0.5),
            'Genesis = 微型游戏引擎 + AI个性化 + 鸿蒙生态  →  打造"零等待"游戏体验',
            font_size=13, color=TEAL, alignment=PP_ALIGN.CENTER)

print("✓ Slide 4: 产品定位与愿景")

# SLIDE 5: 核心价值主张
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 5)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "核心价值主张", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

values = [
    ("⚡", "零安装体验", "基于鸿蒙元服务架构\n无需下载安装包\n点击即玩，退出即走", TEAL),
    ("🧠", "AI智能驱动", "6维用户画像系统\n实时动态难度调节\n个性化游戏推荐", ACCENT_RED),
    ("🎯", "场景触发", "9大智能触发器\n时间/行为/情绪感知\n精准推送合适游戏", GOLD),
    ("🔄", "持续进化", "20局学习曲线\n画像实时进化\n越玩越懂你", TEAL),
    ("📊", "数据闭环", "行为采集→画像更新\n→推荐优化→体验提升\n全链路数据驱动", ACCENT_RED),
    ("🌐", "生态整合", "华为游戏中心入口\n服务卡片桌面直达\n分布式跨设备体验", GOLD),
]

for i, (icon, title, desc, accent) in enumerate(values):
    col = i % 3
    row = i // 3
    x = Inches(0.6) + Inches(3.1) * col
    y = Inches(1.2) + Inches(2.6) * row
    
    card = add_shape(slide, x, y, Inches(2.8), Inches(2.3),
                     fill_color=DARK_PANEL, line_color=accent, line_width=Pt(1))
    add_textbox(slide, x + Inches(0.15), y + Inches(0.1), Inches(0.5), Inches(0.4),
                icon, font_size=20, color=accent)
    add_textbox(slide, x + Inches(0.6), y + Inches(0.15), Inches(2.0), Inches(0.35),
                title, font_size=14, color=accent, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.65), Inches(2.4), Inches(1.4),
                desc, font_size=11, color=LIGHT_GRAY)

print("✓ Slide 5: 核心价值主张")

# SLIDE 6: 目标用户与场景
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 6)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "目标用户与场景", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Target user
add_textbox(slide, Inches(0.6), Inches(1.1), Inches(4), Inches(0.4),
            "🎯 核心用户画像", font_size=16, color=GOLD, bold=True)

user_profiles = [
    ("游戏中心活跃用户", "下载/更新游戏频繁，等待时间多"),
    ("休闲游戏爱好者", "碎片化时间，偏好轻量体验"),
    ("鸿蒙生态用户", "HarmonyOS设备，习惯元服务"),
]
for i, (profile, desc) in enumerate(user_profiles):
    y = Inches(1.6) + Inches(0.6) * i
    add_shape(slide, Inches(0.6), y, Inches(0.15), Inches(0.4), fill_color=TEAL, shape_type=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, Inches(0.9), y, Inches(2), Inches(0.3),
                profile, font_size=12, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.9), y + Inches(0.25), Inches(3.5), Inches(0.25),
                desc, font_size=10, color=LIGHT_GRAY)

# Scenarios
add_textbox(slide, Inches(5.2), Inches(1.1), Inches(4), Inches(0.4),
            "📱 典型使用场景", font_size=16, color=GOLD, bold=True)

scenarios = [
    ("⏬ 下载等待", "大游戏下载时推送小游戏消磨时间"),
    ("😔 连败安慰", "连续失败后推送放松类小游戏"),
    ("🌙 深夜关怀", "深夜时段推送助眠呼吸游戏"),
    ("☀️ 早安问候", "早晨推送益智类游戏唤醒大脑"),
    ("🚇 通勤陪伴", "检测移动状态推送快速上手游戏"),
    ("🌧️ 雨天治愈", "天气API检测雨天推送治愈系游戏"),
    ("👋 卸载挽回", "用户卸载游戏时推送挽留小游戏"),
    ("👥 好友在线", "好友在线时推送社交互动小游戏"),
]

for i, (scene, desc) in enumerate(scenarios):
    col = i % 2
    row = i // 2
    x = Inches(5.0) + Inches(2.3) * col
    y = Inches(1.6) + Inches(0.65) * row
    add_shape(slide, x, y, Inches(2.2), Inches(0.55), fill_color=DARK_PANEL,
              line_color=RGBColor(0x33, 0x33, 0x55), line_width=Pt(0.5))
    add_textbox(slide, x + Inches(0.1), y + Inches(0.02), Inches(2.0), Inches(0.22),
                scene, font_size=10, color=TEAL, bold=True)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.25), Inches(2.0), Inches(0.25),
                desc, font_size=8, color=LIGHT_GRAY)

# Funnel at bottom
add_textbox(slide, Inches(0.6), Inches(4.8), Inches(8.8), Inches(0.4),
            "用户旅程", font_size=14, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

journey_steps = ["触发场景", "推送卡片", "点击即玩", "AI个性化", "数据反馈", "画像进化"]
for i, step in enumerate(journey_steps):
    x = Inches(0.4) + Inches(1.55) * i
    color = [ACCENT_RED, TEAL, GOLD, TEAL, ACCENT_RED, GOLD][i]
    add_arrow_box(slide, x, Inches(5.3), Inches(1.5), Inches(0.5), step, color=color)

print("✓ Slide 6: 目标用户与场景")

# ────────────────────────────────────────────────────────────
# PART 2: 技术架构
# ────────────────────────────────────────────────────────────

# SLIDE 7: Section Divider - Part 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "02", "技术架构", "Technical Architecture")
add_slide_number(slide, 7)
print("✓ Slide 7: Part 2 Section")

# SLIDE 8: 系统架构图
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 8)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "系统架构图  ·  分层架构", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Architecture layers
layers = [
    ("HarmonyOS 集成层", "EntryAbility | 服务卡片 | ArkTS Bridge | 分布式能力", ACCENT_RED, Inches(1.1)),
    ("触发器层 Triggers", "下载等待 | 连败安慰 | 深夜关怀 | 早安问候 | 通勤陪伴 | 无聊检测 | 卸载挽回 | 雨天治愈 | 好友在线", RGBColor(0xBB, 0x33, 0x33), Inches(1.9)),
    ("模板层 Templates", "消除 | 节奏 | 翻牌 | 呼吸 | 绘画 | 颜色 | 算术 | 泡泡 | 单词 | 追光 | 2048 | 贪吃蛇 | 弹球 | 打地鼠 | 接金币", RGBColor(0x88, 0x66, 0x22), Inches(2.7)),
    ("核心引擎层 Engine", "MicroGameEngine | TriggerEngine | PersonalizationEngine | TemplateFactory | DifficultyManager", TEAL, Inches(3.5)),
    ("AI 模型层 Models", "UserProfile | TasteProfile | DifficultyCurve | SceneContext | RecommendationEngine | LearningTracker", GOLD, Inches(4.3)),
    ("数据层 Data", "TypeScript 5.9.3 | Event System | State Management | Analytics | Storage", RGBColor(0x33, 0x66, 0x99), Inches(5.1)),
]

for name, desc, color, y_pos in layers:
    # Layer background
    add_shape(slide, Inches(0.5), y_pos, Inches(9.0), Inches(0.7),
              fill_color=DARK_PANEL, line_color=color, line_width=Pt(1.5))
    # Color bar on left
    add_shape(slide, Inches(0.5), y_pos, Inches(0.08), Inches(0.7),
              fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, Inches(0.7), y_pos + Inches(0.05), Inches(2.5), Inches(0.3),
                name, font_size=12, color=color, bold=True)
    add_textbox(slide, Inches(0.7), y_pos + Inches(0.32), Inches(8.5), Inches(0.3),
                desc, font_size=9, color=LIGHT_GRAY)

# Arrows between layers
for i in range(5):
    y = Inches(1.85) + Inches(0.8) * i
    add_textbox(slide, Inches(4.5), y, Inches(1), Inches(0.25),
                "▲ ▼", font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Side label
add_textbox(slide, Inches(9.2), Inches(1.5), Inches(0.7), Inches(4),
            "承\n上\n启\n下", font_size=12, color=MEDIUM_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

print("✓ Slide 8: 系统架构图")

# SLIDE 9: 核心引擎模块
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 9)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "核心引擎模块", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

modules = [
    ("MicroGameEngine", "游戏引擎核心", [
        "统一游戏生命周期管理",
        "start() → play() → pause() → end()",
        "事件驱动的状态机",
        "分数/时间/进度统一接口",
    ], ACCENT_RED),
    ("TriggerEngine", "智能触发引擎", [
        "9种场景触发器管理",
        "多信号融合决策",
        "优先级排序与冲突解决",
        "触发频率控制与冷却",
    ], TEAL),
    ("PersonalizationEngine", "个性化引擎", [
        "6维用户画像构建",
        "实时推荐算法",
        "动态难度调节",
        "学习曲线追踪",
    ], GOLD),
    ("TemplateFactory", "模板工厂", [
        "15种游戏模板注册",
        "配置驱动的模板实例化",
        "主题/参数动态注入",
        "模板版本管理",
    ], RGBColor(0xAA, 0x55, 0xFF)),
]

for i, (name, subtitle, features, color) in enumerate(modules):
    x = Inches(0.4) + Inches(2.35) * i
    # Module card
    add_shape(slide, x, Inches(1.2), Inches(2.2), Inches(5.2),
              fill_color=DARK_PANEL, line_color=color, line_width=Pt(1.5))
    # Header
    add_shape(slide, x, Inches(1.2), Inches(2.2), Inches(0.9), fill_color=color,
              shape_type=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, x + Inches(0.1), Inches(1.25), Inches(2.0), Inches(0.4),
                name, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), Inches(1.65), Inches(2.0), Inches(0.35),
                subtitle, font_size=10, color=RGBColor(0xFF, 0xFF, 0xDD), alignment=PP_ALIGN.CENTER)
    
    # Features
    for j, feat in enumerate(features):
        y = Inches(2.3) + Inches(0.55) * j
        add_shape(slide, x + Inches(0.15), y, Inches(0.12), Inches(0.12),
                  fill_color=color, shape_type=MSO_SHAPE.OVAL)
        add_textbox(slide, x + Inches(0.35), y - Inches(0.05), Inches(1.7), Inches(0.5),
                    feat, font_size=10, color=LIGHT_GRAY)

print("✓ Slide 9: 核心引擎模块")

# SLIDE 10: 15个游戏模板矩阵
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 10)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "15个游戏模板矩阵", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Table data for 15 games
game_data = [
    ["#", "游戏名称", "类型", "难度", "平均时长", "触发场景"],
    ["1", "消除方块", "益智", "★☆☆☆", "30s", "下载等待"],
    ["2", "节奏大师", "音乐", "★★☆☆", "45s", "无聊检测"],
    ["3", "翻牌配对", "记忆", "★☆☆☆", "25s", "早安问候"],
    ["4", "呼吸冥想", "放松", "☆☆☆☆", "60s", "深夜/连败"],
    ["5", "绘画挑战", "创意", "★★☆☆", "40s", "通勤陪伴"],
    ["6", "颜色识别", "反应", "★☆☆☆", "20s", "下载等待"],
    ["7", "速算挑战", "益智", "★★★☆", "30s", "早安问候"],
    ["8", "泡泡射手", "休闲", "★☆☆☆", "35s", "无聊检测"],
    ["9", "单词拼图", "教育", "★★☆☆", "40s", "早安问候"],
    ["10", "追光游戏", "反应", "★★☆☆", "25s", "下载等待"],
    ["11", "2048", "益智", "★★★☆", "60s", "通勤陪伴"],
    ["12", "贪吃蛇", "经典", "★★☆☆", "45s", "无聊检测"],
    ["13", "弹球大师", "动作", "★★★☆", "35s", "好友在线"],
    ["14", "打地鼠", "反应", "★★☆☆", "30s", "雨天治愈"],
    ["15", "接金币", "休闲", "★☆☆☆", "25s", "卸载挽回"],
]

tbl = add_table(slide, Inches(0.4), Inches(1.1), Inches(9.2), Inches(5.6),
                len(game_data), 6, game_data,
                col_widths=[Inches(0.5), Inches(1.3), Inches(1.0), Inches(1.2), Inches(1.2), Inches(1.5)])

print("✓ Slide 10: 15个游戏模板矩阵")

# SLIDE 11: 9大智能触发器
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 11)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "9大智能触发器", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

triggers = [
    ("⏬", "下载等待", "检测游戏下载/更新\n推送快速上手小游戏", "高", ACCENT_RED),
    ("😔", "连败安慰", "连续3+局失败\n推送放松治愈游戏", "高", ACCENT_RED),
    ("🌙", "深夜关怀", "22:00-06:00时段\n推送助眠呼吸游戏", "中", RGBColor(0x66, 0x66, 0xCC)),
    ("☀️", "早安问候", "06:00-09:00时段\n推送益智唤醒游戏", "中", GOLD),
    ("🚇", "通勤陪伴", "加速度计+时间窗口\n推送中度游戏", "中", TEAL),
    ("😴", "无聊检测", "低活跃+重复操作\n推送新鲜游戏", "高", ACCENT_RED),
    ("👋", "卸载挽回", "卸载确认时触发\n推送经典小游戏", "极高", RGBColor(0xFF, 0x66, 0x33)),
    ("🌧️", "雨天治愈", "天气API+地理位置\n推送治愈系游戏", "低", RGBColor(0x66, 0x99, 0xCC)),
    ("👥", "好友在线", "社交信号+好友状态\n推送互动小游戏", "中", TEAL),
]

for i, (icon, name, desc, priority, color) in enumerate(triggers):
    col = i % 3
    row = i // 3
    x = Inches(0.4) + Inches(3.15) * col
    y = Inches(1.15) + Inches(2.0) * row
    
    add_shape(slide, x, y, Inches(2.95), Inches(1.8),
              fill_color=DARK_PANEL, line_color=color, line_width=Pt(1))
    # Icon and name
    add_textbox(slide, x + Inches(0.1), y + Inches(0.1), Inches(0.4), Inches(0.35),
                icon, font_size=18, color=color)
    add_textbox(slide, x + Inches(0.5), y + Inches(0.1), Inches(1.5), Inches(0.35),
                name, font_size=14, color=WHITE, bold=True)
    # Priority badge
    badge_color = ACCENT_RED if priority in ["高", "极高"] else TEAL
    add_shape(slide, x + Inches(2.1), y + Inches(0.1), Inches(0.7), Inches(0.3),
              fill_color=badge_color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, x + Inches(2.1), y + Inches(0.1), Inches(0.7), Inches(0.3),
                priority, font_size=8, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_textbox(slide, x + Inches(0.15), y + Inches(0.6), Inches(2.6), Inches(1.0),
                desc, font_size=10, color=LIGHT_GRAY)

print("✓ Slide 11: 9大智能触发器")

# SLIDE 12: 数据流图
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 12)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "数据流图  ·  全链路数据闭环", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Flow stages
flow_stages = [
    ("1. 用户行为", "点击/滑动/停留\n游戏选择/时长/分数", ACCENT_RED),
    ("2. 触发评估", "场景信号采集\n多条件匹配决策", TEAL),
    ("3. 画像匹配", "6维偏好匹配\n置信度评分排序", GOLD),
    ("4. 游戏生成", "模板实例化\n参数动态注入", RGBColor(0xAA, 0x55, 0xFF)),
    ("5. 结果反馈", "游戏数据采集\n体验评分记录", TEAL),
    ("6. 画像更新", "偏好权重调整\n难度曲线学习", ACCENT_RED),
]

for i, (title, desc, color) in enumerate(flow_stages):
    x = Inches(0.3) + Inches(1.6) * i
    y = Inches(1.2)
    
    # Box
    add_shape(slide, x, y, Inches(1.45), Inches(1.8),
              fill_color=DARK_PANEL, line_color=color, line_width=Pt(1.5))
    add_shape(slide, x, y, Inches(1.45), Inches(0.35), fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, x + Inches(0.05), y + Inches(0.02), Inches(1.35), Inches(0.3),
                title, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.5), Inches(1.25), Inches(1.2),
                desc, font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    
    # Arrow
    if i < 5:
        add_textbox(slide, x + Inches(1.35), y + Inches(0.6), Inches(0.4), Inches(0.3),
                    "→", font_size=16, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Cycle arrow at bottom
add_shape(slide, Inches(0.3), Inches(3.5), Inches(9.4), Pt(2), fill_color=MEDIUM_GRAY, shape_type=MSO_SHAPE.RECTANGLE)
add_textbox(slide, Inches(3.5), Inches(3.65), Inches(3), Inches(0.3),
            "⟳  持续循环优化", font_size=12, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

# Detail section
add_textbox(slide, Inches(0.6), Inches(4.2), Inches(4), Inches(0.35),
            "关键数据指标", font_size=16, color=GOLD, bold=True)

metrics = [
    ("画像维度", "6维 (策略/动作/社交/创意/放松/竞技)"),
    ("推荐延迟", "< 50ms 实时响应"),
    ("学习周期", "20局达到画像收敛"),
    ("数据采集", "50+ 行为维度实时记录"),
    ("隐私保护", "本地计算为主，脱敏上传"),
]

for i, (label, value) in enumerate(metrics):
    y = Inches(4.7) + Inches(0.4) * i
    add_textbox(slide, Inches(0.8), y, Inches(1.5), Inches(0.3),
                label, font_size=11, color=TEAL, bold=True)
    add_textbox(slide, Inches(2.5), y, Inches(3), Inches(0.3),
                value, font_size=11, color=LIGHT_GRAY)

# Tech stack
add_textbox(slide, Inches(5.5), Inches(4.2), Inches(4), Inches(0.35),
            "技术栈", font_size=16, color=GOLD, bold=True)

stack_items = [
    "TypeScript 5.9.3  核心引擎语言",
    "Event System  事件驱动架构",
    "State Machine  游戏状态管理",
    "Observer Pattern  画像更新通知",
    "Strategy Pattern  触发器策略",
]
for i, item in enumerate(stack_items):
    y = Inches(4.7) + Inches(0.4) * i
    add_textbox(slide, Inches(5.7), y, Inches(4), Inches(0.3),
                f"▸ {item}", font_size=11, color=LIGHT_GRAY)

print("✓ Slide 12: 数据流图")

# ────────────────────────────────────────────────────────────
# PART 3: AI能力
# ────────────────────────────────────────────────────────────

# SLIDE 13: Section Divider - Part 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "03", "AI 能力", "AI Capabilities")
add_slide_number(slide, 13)
print("✓ Slide 13: Part 3 Section")

# SLIDE 14: 用户画像系统
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 14)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "用户画像系统  ·  6维味道雷达图", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Radar chart simulation (hexagonal representation)
center_x = Inches(2.8)
center_y = Inches(3.5)
radius = Inches(1.5)

dimensions = [
    ("策略", "STRATEGY", ACCENT_RED, "偏好思考规划类\n如2048、消除"),
    ("动作", "ACTION", TEAL, "偏好反应操作类\n如打地鼠、追光"),
    ("社交", "SOCIAL", GOLD, "偏好互动分享类\n如好友对战"),
    ("创意", "CREATIVE", RGBColor(0xAA, 0x55, 0xFF), "偏好自由表达类\n如绘画挑战"),
    ("放松", "RELAXATION", RGBColor(0x66, 0xCC, 0x66), "偏好舒缓减压类\n如呼吸冥想"),
    ("竞技", "COMPETITIVE", RGBColor(0xFF, 0x88, 0x33), "偏好排名挑战类\n如速算、弹球"),
]

for i, (name, en, color, desc) in enumerate(dimensions):
    angle_offset = i * 60
    # Position labels around center
    import math
    angle = math.radians(-90 + angle_offset)
    lx = Inches(2.8) + Inches(1.8) * math.cos(angle)
    ly = Inches(3.3) + Inches(1.8) * math.sin(angle)
    
    # Dimension label
    add_textbox(slide, lx - Inches(0.6), ly - Inches(0.15), Inches(1.2), Inches(0.35),
                f"{name}", font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, lx - Inches(0.6), ly + Inches(0.15), Inches(1.2), Inches(0.2),
                en, font_size=7, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Center hub
add_shape(slide, Inches(2.4), Inches(3.0), Inches(0.8), Inches(0.8),
          fill_color=ACCENT_RED, line_color=GOLD, line_width=Pt(2), shape_type=MSO_SHAPE.OVAL)
add_textbox(slide, Inches(2.4), Inches(3.15), Inches(0.8), Inches(0.5),
            "Taste\nProfile", font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Right side: Taste → GameType mapping
add_textbox(slide, Inches(5.2), Inches(1.1), Inches(4.5), Inches(0.4),
            "Taste → 游戏类型映射", font_size=16, color=GOLD, bold=True)

mapping_data = [
    ["Taste维度", "推荐游戏类型", "置信度"],
    ["策略 ↑↑↑", "2048 / 消除 / 单词", "92%"],
    ["动作 ↑↑", "打地鼠 / 追光 / 弹球", "87%"],
    ["社交 ↑↑", "好友对战 / 接金币", "78%"],
    ["创意 ↑", "绘画 / 颜色识别", "85%"],
    ["放松 ↑↑↑", "呼吸 / 泡泡 / 消除", "95%"],
    ["竞技 ↑↑", "速算 / 2048 / 弹球", "89%"],
]

add_table(slide, Inches(5.2), Inches(1.6), Inches(4.4), Inches(3.2),
          len(mapping_data), 3, mapping_data,
          col_widths=[Inches(1.3), Inches(1.8), Inches(1.0)])

# Bottom note
add_shape(slide, Inches(0.5), Inches(5.5), Inches(9.0), Inches(1.2), fill_color=DARK_PANEL,
          line_color=TEAL, line_width=Pt(1))
add_textbox(slide, Inches(0.7), Inches(5.6), Inches(8.5), Inches(0.3),
            "📊 画像更新机制", font_size=13, color=TEAL, bold=True)
add_textbox(slide, Inches(0.7), Inches(5.95), Inches(8.5), Inches(0.6),
            "每局游戏结束后自动更新 | 基于指数移动平均(EMA)平滑调整 | 20局后达到稳定收敛 | 本地存储保护隐私",
            font_size=11, color=LIGHT_GRAY)

print("✓ Slide 14: 用户画像系统")

# SLIDE 15: 智能推荐引擎
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 15)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "智能推荐引擎", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Recommendation flow
add_textbox(slide, Inches(0.6), Inches(1.1), Inches(4), Inches(0.35),
            "推荐算法流程", font_size=16, color=GOLD, bold=True)

rec_steps = [
    ("输入信号", "用户画像 + 场景上下文 + 历史记录", ACCENT_RED),
    ("候选生成", "Taste→GameType映射生成候选集", TEAL),
    ("特征提取", "难度匹配度 + 时长适配 + 新鲜度 + 历史偏好", GOLD),
    ("评分排序", "加权评分 = Σ(wi × fi)  综合排序", RGBColor(0xAA, 0x55, 0xFF)),
    ("输出推荐", "Top-3推荐 + 置信度 + 推荐理由", TEAL),
]

for i, (step, desc, color) in enumerate(rec_steps):
    y = Inches(1.6) + Inches(0.9) * i
    add_shape(slide, Inches(0.6), y, Inches(4.2), Inches(0.75),
              fill_color=DARK_PANEL, line_color=color, line_width=Pt(1))
    add_shape(slide, Inches(0.6), y, Inches(0.08), Inches(0.75),
              fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, Inches(0.85), y + Inches(0.05), Inches(1.2), Inches(0.3),
                step, font_size=12, color=color, bold=True)
    add_textbox(slide, Inches(0.85), y + Inches(0.35), Inches(3.8), Inches(0.35),
                desc, font_size=10, color=LIGHT_GRAY)
    if i < 4:
        add_textbox(slide, Inches(2.3), y + Inches(0.7), Inches(0.5), Inches(0.25),
                    "↓", font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Right: Scoring formula
add_textbox(slide, Inches(5.2), Inches(1.1), Inches(4.5), Inches(0.35),
            "置信度评分模型", font_size=16, color=GOLD, bold=True)

formula_items = [
    ("难度匹配权重", "0.30", "基于用户历史表现"),
    ("时长适配权重", "0.25", "匹配场景可用时间"),
    ("新鲜度权重", "0.20", "避免重复推荐"),
    ("历史偏好权重", "0.15", "强化正面体验"),
    ("多样性权重", "0.10", "探索新类型"),
]

for i, (name, weight, desc) in enumerate(formula_items):
    y = Inches(1.6) + Inches(0.75) * i
    add_shape(slide, Inches(5.2), y, Inches(4.4), Inches(0.6),
              fill_color=DARK_PANEL, line_color=RGBColor(0x33, 0x33, 0x55), line_width=Pt(0.5))
    
    # Weight badge
    add_shape(slide, Inches(5.3), y + Inches(0.12), Inches(0.8), Inches(0.35),
              fill_color=TEAL, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(5.3), y + Inches(0.12), Inches(0.8), Inches(0.35),
                weight, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(6.3), y + Inches(0.08), Inches(3.2), Inches(0.25),
                name, font_size=11, color=WHITE, bold=True)
    add_textbox(slide, Inches(6.3), y + Inches(0.32), Inches(3.2), Inches(0.25),
                desc, font_size=9, color=LIGHT_GRAY)

# Bottom example
add_shape(slide, Inches(0.5), Inches(5.8), Inches(9.0), Inches(1.0), fill_color=DARK_PANEL,
          line_color=ACCENT_RED, line_width=Pt(1))
add_textbox(slide, Inches(0.7), Inches(5.9), Inches(8.5), Inches(0.3),
            "💡 推荐示例", font_size=13, color=ACCENT_RED, bold=True)
add_textbox(slide, Inches(0.7), Inches(6.25), Inches(8.5), Inches(0.4),
            "用户画像: 策略↑↑ + 放松↑  |  场景: 深夜23:00  |  →  推荐消除方块 (92%置信度) + 呼吸冥想 (85%置信度)",
            font_size=11, color=LIGHT_GRAY)

print("✓ Slide 15: 智能推荐引擎")

# SLIDE 16: 动态难度调节
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 16)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "动态难度调节  ·  基于表现的实时调整", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Difficulty levels
diff_levels = [
    ("EASY", "简单", "初学者", "0-30%", RGBColor(0x66, 0xCC, 0x66), "3x3网格 / 慢速 / 提示多"),
    ("NORMAL", "普通", "进阶者", "30-60%", TEAL, "4x4网格 / 中速 / 适度提示"),
    ("HARD", "困难", "高手", "60-85%", GOLD, "5x5网格 / 快速 / 少提示"),
    ("EXPERT", "专家", "大师", "85-100%", ACCENT_RED, "6x6网格 / 极速 / 无提示"),
]

for i, (level, cn, role, range_val, color, desc) in enumerate(diff_levels):
    x = Inches(0.5) + Inches(2.35) * i
    
    # Level card
    add_shape(slide, x, Inches(1.2), Inches(2.2), Inches(2.5),
              fill_color=DARK_PANEL, line_color=color, line_width=Pt(2))
    
    # Level header
    add_shape(slide, x, Inches(1.2), Inches(2.2), Inches(0.7), fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, x, Inches(1.22), Inches(2.2), Inches(0.35),
                level, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(1.55), Inches(2.2), Inches(0.3),
                f"{cn} · {role}", font_size=10, color=RGBColor(0xFF, 0xFF, 0xDD), alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, x + Inches(0.15), Inches(2.05), Inches(1.9), Inches(0.3),
                f"得分区间: {range_val}", font_size=10, color=color, bold=True)
    add_textbox(slide, x + Inches(0.15), Inches(2.4), Inches(1.9), Inches(1.0),
                desc, font_size=10, color=LIGHT_GRAY)

# Transition arrows
for i in range(3):
    x = Inches(2.55) + Inches(2.35) * i
    add_textbox(slide, x, Inches(2.2), Inches(0.4), Inches(0.3),
                "→", font_size=18, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Adjustment logic
add_textbox(slide, Inches(0.6), Inches(4.0), Inches(8), Inches(0.35),
            "📈 调节算法", font_size=16, color=GOLD, bold=True)

adj_data = [
    ["条件", "动作", "幅度"],
    ["连续3局得分 > 80%", "升级难度", "+1级"],
    ["连续3局得分 < 30%", "降低难度", "-1级"],
    ["单局得分 > 95%", "标记潜力", "加速升级"],
    ["连续5局波动大", "维持当前", "稳定观察"],
    ["新用户前3局", "从EASY开始", "快速校准"],
]

add_table(slide, Inches(0.5), Inches(4.5), Inches(9.0), Inches(2.3),
          len(adj_data), 3, adj_data,
          col_widths=[Inches(3.0), Inches(3.0), Inches(2.0)])

print("✓ Slide 16: 动态难度调节")

# SLIDE 17: 学习曲线
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 17)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "学习曲线  ·  20局画像进化可视化", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Simulated learning curve visualization
add_textbox(slide, Inches(0.6), Inches(1.1), Inches(8.8), Inches(0.35),
            "用户画像收敛过程 (6维度随游戏局数变化)", font_size=14, color=GOLD, bold=True)

# Y-axis
add_shape(slide, Inches(1.2), Inches(1.7), Pt(2), Inches(2.5), fill_color=MEDIUM_GRAY, shape_type=MSO_SHAPE.RECTANGLE)
add_textbox(slide, Inches(0.5), Inches(1.7), Inches(0.6), Inches(2.5),
            "100%\n\n75%\n\n50%\n\n25%\n\n0%", font_size=8, color=MEDIUM_GRAY, alignment=PP_ALIGN.RIGHT)

# X-axis
add_shape(slide, Inches(1.2), Inches(4.2), Inches(7.8), Pt(2), fill_color=MEDIUM_GRAY, shape_type=MSO_SHAPE.RECTANGLE)
add_textbox(slide, Inches(1.2), Inches(4.3), Inches(7.8), Inches(0.25),
            "第1局    第3局    第5局    第8局    第10局    第13局    第15局    第18局    第20局",
            font_size=8, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Phase markers
phases = [
    ("探索期\n(1-5局)", "全维度尝试\n画像波动大", ACCENT_RED, Inches(1.2)),
    ("校准期\n(6-12局)", "偏好初现\n收敛趋势", GOLD, Inches(3.8)),
    ("稳定期\n(13-20局)", "画像收敛\n个性化精准", RGBColor(0x66, 0xCC, 0x66), Inches(6.4)),
]

for label, desc, color, x in phases:
    # Phase region
    add_shape(slide, x, Inches(1.7), Inches(2.5), Inches(2.5),
              fill_color=None, line_color=color, line_width=Pt(1), shape_type=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, x, Inches(1.75), Inches(2.5), Inches(0.5),
                label, font_size=9, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(2.3), Inches(2.5), Inches(0.5),
                desc, font_size=8, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Key insights
add_textbox(slide, Inches(0.6), Inches(4.8), Inches(8.8), Inches(0.35),
            "🔍 关键发现", font_size=16, color=GOLD, bold=True)

insights = [
    ("前5局", "推荐多样化游戏类型，广泛收集偏好数据，画像置信度 < 50%", ACCENT_RED),
    ("6-12局", "核心偏好浮现，推荐精准度提升至70%+，开始个性化调整", TEAL),
    ("13-20局", "画像稳定收敛，推荐命中率 > 85%，用户体验显著提升", RGBColor(0x66, 0xCC, 0x66)),
    ("20局后", "进入精调阶段，微调权重，引入新奇探索避免过度拟合", GOLD),
]

for i, (phase, insight, color) in enumerate(insights):
    y = Inches(5.3) + Inches(0.45) * i
    add_shape(slide, Inches(0.8), y + Inches(0.05), Inches(0.8), Inches(0.3),
              fill_color=color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(0.8), y + Inches(0.05), Inches(0.8), Inches(0.3),
                phase, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.8), y, Inches(7.5), Inches(0.35),
                insight, font_size=10, color=LIGHT_GRAY)

print("✓ Slide 17: 学习曲线")

# SLIDE 18: 场景感知触发
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 18)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "场景感知触发  ·  多信号融合决策", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Signal types
signals = [
    ("⏰", "时间信号", "时段 / 星期 / 节假日", "早晚通勤、深夜、午休", ACCENT_RED),
    ("🎮", "行为信号", "游戏记录 / 点击流 / 停留", "下载中、连败、低活跃", TEAL),
    ("😊", "情绪信号", "面部表情 / 输入速度 / 操作", "挫败、无聊、兴奋", GOLD),
    ("🌤️", "环境信号", "天气 / 位置 / 温度", "下雨天、周末在家", RGBColor(0xAA, 0x55, 0xFF)),
    ("👥", "社交信号", "好友状态 / 分享记录", "好友在线、好友分享", RGBColor(0x66, 0xCC, 0x66)),
]

for i, (icon, name, sources, examples, color) in enumerate(signals):
    y = Inches(1.2) + Inches(0.95) * i
    
    add_shape(slide, Inches(0.5), y, Inches(4.2), Inches(0.8),
              fill_color=DARK_PANEL, line_color=color, line_width=Pt(1))
    add_textbox(slide, Inches(0.6), y + Inches(0.1), Inches(0.4), Inches(0.35),
                icon, font_size=18, color=color)
    add_textbox(slide, Inches(1.05), y + Inches(0.05), Inches(1.5), Inches(0.3),
                name, font_size=13, color=color, bold=True)
    add_textbox(slide, Inches(1.05), y + Inches(0.35), Inches(3.5), Inches(0.3),
                f"采集: {sources}", font_size=9, color=LIGHT_GRAY)
    add_textbox(slide, Inches(2.7), y + Inches(0.05), Inches(2.0), Inches(0.3),
                f"示例: {examples}", font_size=9, color=MEDIUM_GRAY)

# Fusion engine
add_textbox(slide, Inches(5.2), Inches(1.1), Inches(4.5), Inches(0.35),
            "多信号融合决策模型", font_size=16, color=GOLD, bold=True)

# Fusion process
fusion_steps = [
    ("信号采集层", "5类信号并行采集", ACCENT_RED),
    ("特征提取层", "标准化 + 归一化", TEAL),
    ("权重分配层", "动态权重 (0.0-1.0)", GOLD),
    ("融合决策层", "加权评分 + 阈值判定", RGBColor(0xAA, 0x55, 0xFF)),
    ("触发输出层", "触发/抑制/排队", RGBColor(0x66, 0xCC, 0x66)),
]

for i, (name, desc, color) in enumerate(fusion_steps):
    y = Inches(1.6) + Inches(0.85) * i
    add_shape(slide, Inches(5.2), y, Inches(4.3), Inches(0.7),
              fill_color=DARK_PANEL, line_color=color, line_width=Pt(1))
    add_shape(slide, Inches(5.2), y, Inches(0.08), Inches(0.7),
              fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, Inches(5.45), y + Inches(0.05), Inches(2.0), Inches(0.3),
                name, font_size=11, color=color, bold=True)
    add_textbox(slide, Inches(5.45), y + Inches(0.32), Inches(3.8), Inches(0.3),
                desc, font_size=9, color=LIGHT_GRAY)
    if i < 4:
        add_textbox(slide, Inches(7.0), y + Inches(0.62), Inches(1), Inches(0.2),
                    "↓", font_size=12, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Example scenario
add_shape(slide, Inches(0.5), Inches(6.0), Inches(9.0), Inches(0.8), fill_color=DARK_PANEL,
          line_color=ACCENT_RED, line_width=Pt(1))
add_textbox(slide, Inches(0.7), Inches(6.05), Inches(8.5), Inches(0.3),
            "💡 融合示例", font_size=12, color=ACCENT_RED, bold=True)
add_textbox(slide, Inches(0.7), Inches(6.35), Inches(8.5), Inches(0.3),
            "时间=23:00(深夜) + 行为=连续3败(挫败) + 环境=雨天(低落) → 触发「深夜关怀」→ 推荐呼吸冥想 (置信度98%)",
            font_size=10, color=LIGHT_GRAY)

print("✓ Slide 18: 场景感知触发")

# ────────────────────────────────────────────────────────────
# PART 4: 鸿蒙集成
# ────────────────────────────────────────────────────────────

# SLIDE 19: Section Divider - Part 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "04", "鸿蒙集成", "HarmonyOS Integration")
add_slide_number(slide, 19)
print("✓ Slide 19: Part 4 Section")

# SLIDE 20: HarmonyOS元服务架构
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 20)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "HarmonyOS 元服务架构", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Three core features
features = [
    ("🚀", "免安装", "元服务无需安装APK\n用户点击即可使用\n冷启动 < 500ms", [
        "Service Widget 技术底座",
        "按需加载游戏资源",
        "缓存策略优化体验",
    ]),
    ("📋", "服务卡片", "2×2 / 2×4 桌面卡片\n实时展示游戏状态\n一键直达游戏", [
        "动态数据刷新",
        "深色/浅色主题适配",
        "交互式卡片操作",
    ]),
    ("🔗", "分布式", "跨设备无缝流转\n手机/平板/手表\n数据云端同步", [
        "HarmonyOS分布式能力",
        "设备发现与迁移",
        "游戏进度实时同步",
    ]),
]

for i, (icon, title, desc, items) in enumerate(features):
    x = Inches(0.4) + Inches(3.2) * i
    
    add_shape(slide, x, Inches(1.2), Inches(3.0), Inches(5.2),
              fill_color=DARK_PANEL, line_color=TEAL, line_width=Pt(1))
    
    # Icon and title
    add_textbox(slide, x + Inches(0.2), Inches(1.35), Inches(0.5), Inches(0.4),
                icon, font_size=24, color=TEAL)
    add_textbox(slide, x + Inches(0.7), Inches(1.35), Inches(2.0), Inches(0.4),
                title, font_size=20, color=GOLD, bold=True)
    
    # Description
    add_textbox(slide, x + Inches(0.2), Inches(1.9), Inches(2.6), Inches(1.2),
                desc, font_size=11, color=LIGHT_GRAY)
    
    # Detail items
    add_accent_line(slide, x + Inches(0.2), Inches(3.2), Inches(2.6), color=TEAL, height=Pt(1))
    for j, item in enumerate(items):
        y = Inches(3.4) + Inches(0.45) * j
        add_textbox(slide, x + Inches(0.2), y, Inches(2.6), Inches(0.4),
                    f"▸ {item}", font_size=10, color=LIGHT_GRAY)

print("✓ Slide 20: HarmonyOS元服务架构")

# SLIDE 21: 服务卡片设计
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 21)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "服务卡片设计", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# 2x2 card mockup
add_textbox(slide, Inches(0.6), Inches(1.1), Inches(4), Inches(0.35),
            "2×2 小卡片", font_size=16, color=GOLD, bold=True)

add_shape(slide, Inches(0.6), Inches(1.6), Inches(2.8), Inches(2.8),
          fill_color=RGBColor(0x2A, 0x2A, 0x4E), line_color=TEAL, line_width=Pt(2))
add_textbox(slide, Inches(0.7), Inches(1.7), Inches(2.6), Inches(0.3),
            "🎮 Genesis", font_size=14, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.7), Inches(2.1), Inches(2.6), Inches(0.3),
            "今日推荐: 消除方块", font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.7), Inches(2.5), Inches(2.6), Inches(0.5),
            "⭐ 2340分 | 🔥 3连胜", font_size=12, color=TEAL, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(1.2), Inches(3.2), Inches(1.6), Inches(0.5),
          fill_color=ACCENT_RED, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, Inches(1.2), Inches(3.2), Inches(1.6), Inches(0.5),
            "点击即玩", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 2x4 card mockup
add_textbox(slide, Inches(5.0), Inches(1.1), Inches(4), Inches(0.35),
            "2×4 宽卡片", font_size=16, color=GOLD, bold=True)

add_shape(slide, Inches(5.0), Inches(1.6), Inches(4.6), Inches(2.8),
          fill_color=RGBColor(0x2A, 0x2A, 0x4E), line_color=GOLD, line_width=Pt(2))
add_textbox(slide, Inches(5.1), Inches(1.7), Inches(4.4), Inches(0.3),
            "🎮 Genesis AI 微游戏", font_size=14, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# Mini game cards inside
mini_games = [("🎯 消除", "推荐"), ("🎵 节奏", "热门"), ("🧘 呼吸", "适合你"), ("🎨 绘画", "新游")]
for i, (game, tag) in enumerate(mini_games):
    x = Inches(5.15) + Inches(1.1) * i
    add_shape(slide, x, Inches(2.2), Inches(1.0), Inches(1.0),
              fill_color=DARK_PANEL, line_color=TEAL, line_width=Pt(0.5))
    add_textbox(slide, x, Inches(2.25), Inches(1.0), Inches(0.5),
                game, font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(2.85), Inches(1.0), Inches(0.25),
                tag, font_size=8, color=TEAL, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(5.1), Inches(3.4), Inches(4.4), Inches(0.3),
            "画像准确度: 92% | 今日已玩: 3局 | 连胜: 🔥5", font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Card specs
add_textbox(slide, Inches(0.6), Inches(4.8), Inches(8.8), Inches(0.35),
            "卡片技术规格", font_size=16, color=GOLD, bold=True)

card_specs = [
    ["规格", "2×2 小卡片", "2×4 宽卡片"],
    ["尺寸", "4格 (约 180×180dp)", "8格 (约 180×360dp)"],
    ["刷新频率", "每30分钟", "每15分钟 / 实时推送"],
    ["显示内容", "推荐游戏 + 分数", "多游戏推荐 + 画像 + 统计"],
    ["交互", "点击打开游戏", "点击选择不同游戏"],
    ["主题", "深色/浅色自动适配", "深色/浅色自动适配"],
]

add_table(slide, Inches(0.5), Inches(5.2), Inches(9.0), Inches(2.0),
          len(card_specs), 3, card_specs,
          col_widths=[Inches(2.0), Inches(3.5), Inches(3.5)])

print("✓ Slide 21: 服务卡片设计")

# SLIDE 22: ArkTS桥接方案
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 22)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "ArkTS 桥接方案", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Architecture diagram
add_textbox(slide, Inches(0.6), Inches(1.1), Inches(8.8), Inches(0.35),
            "TypeScript引擎 → ArkTS UI层 架构", font_size=14, color=GOLD, bold=True)

# Left: TS Engine
add_shape(slide, Inches(0.5), Inches(1.6), Inches(3.5), Inches(3.0),
          fill_color=DARK_PANEL, line_color=TEAL, line_width=Pt(2))
add_textbox(slide, Inches(0.6), Inches(1.65), Inches(3.3), Inches(0.35),
            "TypeScript 引擎层", font_size=14, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

ts_modules = [
    "MicroGameEngine (核心)",
    "TriggerEngine (触发器)",
    "PersonalizationEngine (AI)",
    "TemplateFactory (模板)",
    "15个游戏模板实现",
    "事件系统 & 状态管理",
]
for i, mod in enumerate(ts_modules):
    add_textbox(slide, Inches(0.8), Inches(2.15) + Inches(0.35) * i, Inches(3.0), Inches(0.3),
                f"▸ {mod}", font_size=10, color=LIGHT_GRAY)

# Bridge in middle
add_shape(slide, Inches(4.2), Inches(2.5), Inches(1.8), Inches(1.2),
          fill_color=ACCENT_RED, line_color=GOLD, line_width=Pt(1))
add_textbox(slide, Inches(4.3), Inches(2.55), Inches(1.6), Inches(0.3),
            "Bridge 层", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(4.3), Inches(2.9), Inches(1.6), Inches(0.6),
            "JS↔ArkTS\n数据序列化\n方法映射", font_size=9, color=RGBColor(0xFF, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)

# Arrows
add_textbox(slide, Inches(3.9), Inches(2.8), Inches(0.5), Inches(0.3),
            "→", font_size=18, color=TEAL, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(5.9), Inches(2.8), Inches(0.5), Inches(0.3),
            "→", font_size=18, color=GOLD, alignment=PP_ALIGN.CENTER)

# Right: ArkTS UI
add_shape(slide, Inches(6.3), Inches(1.6), Inches(3.3), Inches(3.0),
          fill_color=DARK_PANEL, line_color=GOLD, line_width=Pt(2))
add_textbox(slide, Inches(6.4), Inches(1.65), Inches(3.1), Inches(0.35),
            "ArkTS UI 层", font_size=14, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

ark_items = [
    "EntryAbility (入口)",
    "GameWidget (服务卡片)",
    "GameUI (游戏界面)",
    "HarmonyOS API 调用",
    "分布式能力调用",
    "本地数据持久化",
]
for i, item in enumerate(ark_items):
    add_textbox(slide, Inches(6.5), Inches(2.15) + Inches(0.35) * i, Inches(3.0), Inches(0.3),
                f"▸ {item}", font_size=10, color=LIGHT_GRAY)

# Key design decisions
add_textbox(slide, Inches(0.6), Inches(5.0), Inches(8.8), Inches(0.35),
            "关键设计决策", font_size=16, color=GOLD, bold=True)

decisions = [
    ("引擎复用", "TypeScript引擎代码100%复用，无需重写"),
    ("UI原生", "ArkTS实现原生UI，性能优于WebView"),
    ("数据桥接", "JSON序列化 + 类型映射，零损耗通信"),
    ("渐进迁移", "可先WebView方案，后续迁移至ArkTS"),
]

for i, (title, desc) in enumerate(decisions):
    col = i % 2
    row = i // 2
    x = Inches(0.5) + Inches(4.7) * col
    y = Inches(5.5) + Inches(0.7) * row
    add_shape(slide, x, y, Inches(4.5), Inches(0.55),
              fill_color=DARK_PANEL, line_color=RGBColor(0x33, 0x33, 0x55), line_width=Pt(0.5))
    add_textbox(slide, x + Inches(0.15), y + Inches(0.05), Inches(1.2), Inches(0.25),
                title, font_size=10, color=TEAL, bold=True)
    add_textbox(slide, x + Inches(1.4), y + Inches(0.05), Inches(2.9), Inches(0.25),
                desc, font_size=10, color=LIGHT_GRAY)

print("✓ Slide 22: ArkTS桥接方案")

# SLIDE 23: EntryAbility生命周期
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 23)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "EntryAbility 生命周期", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Lifecycle flow
lifecycle = [
    ("onCreate()", "初始化引擎\n加载用户画像\n注册触发器", ACCENT_RED, "初始化"),
    ("onForeground()", "检查触发条件\n生成推荐列表\n展示服务卡片", TEAL, "激活"),
    ("onGameStart()", "实例化游戏模板\n注入个性化参数\n启动游戏循环", GOLD, "游戏"),
    ("onBackground()", "暂停游戏\n保存进度\n更新画像", RGBColor(0xAA, 0x55, 0xFF), "暂停"),
    ("onDestroy()", "持久化画像数据\n上报分析数据\n释放资源", ACCENT_RED, "销毁"),
]

for i, (method, desc, color, phase) in enumerate(lifecycle):
    x = Inches(0.3) + Inches(1.95) * i
    
    add_shape(slide, x, Inches(1.2), Inches(1.8), Inches(3.2),
              fill_color=DARK_PANEL, line_color=color, line_width=Pt(1.5))
    
    # Method name
    add_shape(slide, x, Inches(1.2), Inches(1.8), Inches(0.5), fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, x, Inches(1.22), Inches(1.8), Inches(0.45),
                method, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Phase
    add_textbox(slide, x, Inches(1.85), Inches(1.8), Inches(0.3),
                f"[{phase}]", font_size=10, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Description
    add_textbox(slide, x + Inches(0.1), Inches(2.3), Inches(1.6), Inches(1.8),
                desc, font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    
    # Arrow
    if i < 4:
        add_textbox(slide, x + Inches(1.7), Inches(2.3), Inches(0.4), Inches(0.3),
                    "→", font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Timeline at bottom
add_shape(slide, Inches(0.3), Inches(4.7), Inches(9.4), Pt(2), fill_color=MEDIUM_GRAY, shape_type=MSO_SHAPE.RECTANGLE)

time_points = [
    (Inches(0.5), "0ms", "冷启动"),
    (Inches(2.3), "~200ms", "引擎就绪"),
    (Inches(4.3), "~350ms", "推荐完成"),
    (Inches(6.3), "~500ms", "游戏就绪"),
    (Inches(8.3), "<500ms", "可交互"),
]
for x, time_val, label in time_points:
    add_shape(slide, x, Inches(4.6), Pt(4), Pt(12), fill_color=TEAL, shape_type=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, x - Inches(0.3), Inches(4.85), Inches(1), Inches(0.2),
                time_val, font_size=9, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x - Inches(0.3), Inches(5.05), Inches(1), Inches(0.2),
                label, font_size=8, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Key metrics
add_textbox(slide, Inches(0.6), Inches(5.5), Inches(8.8), Inches(0.35),
            "⚡ 性能目标", font_size=16, color=GOLD, bold=True)

perf_metrics = [
    ("冷启动", "< 500ms"),
    ("引擎初始化", "< 200ms"),
    ("AI推荐", "< 50ms"),
    ("游戏加载", "< 100ms"),
    ("画面首帧", "< 150ms"),
]

for i, (metric, target) in enumerate(perf_metrics):
    x = Inches(0.5) + Inches(1.85) * i
    add_shape(slide, x, Inches(6.0), Inches(1.7), Inches(0.8),
              fill_color=DARK_PANEL, line_color=TEAL, line_width=Pt(0.5))
    add_textbox(slide, x, Inches(6.05), Inches(1.7), Inches(0.3),
                metric, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(6.35), Inches(1.7), Inches(0.35),
                target, font_size=14, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

print("✓ Slide 23: EntryAbility生命周期")

# ────────────────────────────────────────────────────────────
# PART 5: 实现成果
# ────────────────────────────────────────────────────────────

# SLIDE 24: Section Divider - Part 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "05", "实现成果", "Implementation Results")
add_slide_number(slide, 24)
print("✓ Slide 24: Part 5 Section")

# SLIDE 25: 工程成果
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 25)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "工程成果", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Key metrics row
key_metrics = [
    ("TypeScript\n5.9.3", "开发语言", TEAL),
    ("15", "游戏模板", GOLD),
    ("363", "测试用例", ACCENT_RED),
    ("94%+", "代码覆盖率", RGBColor(0x66, 0xCC, 0x66)),
]

for i, (val, label, color) in enumerate(key_metrics):
    x = Inches(0.5) + Inches(2.4) * i
    add_shape(slide, x, Inches(1.1), Inches(2.2), Inches(1.4),
              fill_color=DARK_PANEL, line_color=color, line_width=Pt(1.5))
    add_textbox(slide, x, Inches(1.15), Inches(2.2), Inches(0.7),
                val, font_size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(1.9), Inches(2.2), Inches(0.3),
                label, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Code structure
add_textbox(slide, Inches(0.6), Inches(2.8), Inches(4), Inches(0.35),
            "📁 项目结构", font_size=16, color=GOLD, bold=True)

structure_items = [
    "src/engine/         核心引擎 (MicroGameEngine)",
    "src/triggers/       9大触发器实现",
    "src/ai/             AI引擎 (画像/推荐/难度)",
    "src/templates/      15个游戏模板",
    "src/models/         数据模型定义",
    "src/types/          TypeScript类型系统",
    "__tests__/          363个测试文件",
    "demo/               浏览器交互Demo",
]

for i, item in enumerate(structure_items):
    y = Inches(3.3) + Inches(0.35) * i
    add_textbox(slide, Inches(0.8), y, Inches(4.5), Inches(0.3),
                item, font_size=10, color=LIGHT_GRAY, font_name='Consolas')

# Quality metrics
add_textbox(slide, Inches(5.5), Inches(2.8), Inches(4), Inches(0.35),
            "🏆 质量指标", font_size=16, color=GOLD, bold=True)

quality_data = [
    ["指标", "目标", "实际"],
    ["语句覆盖率", "> 90%", "94.2%"],
    ["分支覆盖率", "> 85%", "89.7%"],
    ["函数覆盖率", "> 90%", "96.1%"],
    ["行覆盖率", "> 90%", "94.8%"],
    ["测试套件", "> 300", "363"],
    ["TypeScript严格", "启用", "✓"],
    ["Lint通过率", "100%", "100%"],
]

add_table(slide, Inches(5.3), Inches(3.2), Inches(4.3), Inches(3.0),
          len(quality_data), 3, quality_data,
          col_widths=[Inches(1.6), Inches(1.2), Inches(1.2)])

print("✓ Slide 25: 工程成果")

# SLIDE 26: 浏览器交互Demo
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 26)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "浏览器交互 Demo", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Demo overview
add_textbox(slide, Inches(0.6), Inches(1.1), Inches(8.8), Inches(0.35),
            "10个可玩游戏 + AI Lab 面板  |  浏览器直接体验", font_size=14, color=TEAL, bold=True)

# Playable games
add_textbox(slide, Inches(0.6), Inches(1.6), Inches(4), Inches(0.35),
            "🎮 可玩游戏 (10个)", font_size=16, color=GOLD, bold=True)

playable = [
    "消除方块  —  经典三消玩法",
    "节奏大师  —  音符节拍挑战",
    "翻牌配对  —  记忆力训练",
    "呼吸冥想  —  减压放松",
    "绘画挑战  —  限时绘画",
    "颜色识别  —  反应速度测试",
    "速算挑战  —  数学能力",
    "泡泡射手  —  休闲射击",
    "追光游戏  —  注意力训练",
    "2048      —  策略数字",
]

for i, game in enumerate(playable):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + Inches(2.2) * col
    y = Inches(2.1) + Inches(0.35) * row
    add_textbox(slide, x, y, Inches(2.2), Inches(0.3),
                f"▸ {game}", font_size=9, color=LIGHT_GRAY)

# AI Lab panel
add_textbox(slide, Inches(5.2), Inches(1.6), Inches(4.5), Inches(0.35),
            "🧪 AI Lab 面板", font_size=16, color=GOLD, bold=True)

ai_lab_features = [
    ("用户画像可视化", "实时6维雷达图展示"),
    ("推荐结果调试", "查看推荐算法中间结果"),
    ("触发器模拟器", "模拟不同场景触发"),
    ("难度曲线追踪", "可视化动态难度变化"),
    ("学习进度监控", "20局画像进化实时展示"),
]

for i, (title, desc) in enumerate(ai_lab_features):
    y = Inches(2.1) + Inches(0.55) * i
    add_shape(slide, Inches(5.2), y, Inches(4.4), Inches(0.45),
              fill_color=DARK_PANEL, line_color=RGBColor(0x33, 0x33, 0x55), line_width=Pt(0.5))
    add_textbox(slide, Inches(5.35), y + Inches(0.02), Inches(2), Inches(0.2),
                title, font_size=10, color=TEAL, bold=True)
    add_textbox(slide, Inches(5.35), y + Inches(0.22), Inches(4.0), Inches(0.2),
                desc, font_size=9, color=LIGHT_GRAY)

# Demo access
add_shape(slide, Inches(0.5), Inches(5.5), Inches(9.0), Inches(1.2), fill_color=DARK_PANEL,
          line_color=ACCENT_RED, line_width=Pt(1))
add_textbox(slide, Inches(0.7), Inches(5.6), Inches(8.5), Inches(0.3),
            "🌐 Demo 访问方式", font_size=13, color=ACCENT_RED, bold=True)
add_textbox(slide, Inches(0.7), Inches(5.95), Inches(8.5), Inches(0.6),
            "浏览器打开 index.html 即可体验  |  支持桌面端和移动端  |  内置模拟数据可直接演示AI能力\n推荐使用 Chrome / Edge 最新版  |  移动端建议使用开发者工具的设备模拟模式",
            font_size=11, color=LIGHT_GRAY)

print("✓ Slide 26: 浏览器交互Demo")

# SLIDE 27: CI/CD流水线
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 27)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "CI/CD 流水线", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Pipeline stages
pipeline = [
    ("📥 Code Push", "git push / PR", "代码提交触发", ACCENT_RED),
    ("🔍 Lint & Type", "ESLint + tsc", "代码质量检查", TEAL),
    ("🧪 Test", "Jest 363 tests", "全量测试执行", GOLD),
    ("📊 Coverage", "94%+ coverage", "覆盖率检查", RGBColor(0x66, 0xCC, 0x66)),
    ("🏗️ Build", "tsc + bundle", "生产构建", RGBColor(0xAA, 0x55, 0xFF)),
    ("📦 Release", "GitHub Release", "版本发布", ACCENT_RED),
]

for i, (stage, tool, desc, color) in enumerate(pipeline):
    x = Inches(0.3) + Inches(1.6) * i
    y = Inches(1.2)
    
    add_shape(slide, x, y, Inches(1.5), Inches(2.0),
              fill_color=DARK_PANEL, line_color=color, line_width=Pt(1.5))
    add_textbox(slide, x, y + Inches(0.1), Inches(1.5), Inches(0.35),
                stage, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(0.5), Inches(1.5), Inches(0.3),
                tool, font_size=9, color=color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(0.85), Inches(1.5), Inches(0.8),
                desc, font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    
    if i < 5:
        add_textbox(slide, x + Inches(1.35), y + Inches(0.6), Inches(0.4), Inches(0.3),
                    "→", font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# GitHub Actions details
add_textbox(slide, Inches(0.6), Inches(3.5), Inches(8.8), Inches(0.35),
            "⚙️ GitHub Actions 工作流", font_size=16, color=GOLD, bold=True)

workflows = [
    ["工作流", "触发条件", "执行内容", "预计耗时"],
    ["CI", "push/PR to main", "lint + type + test + coverage", "~3 min"],
    ["Build", "tag push (v*)", "build + bundle + artifact", "~2 min"],
    ["Release", "tag push (v*)", "changelog + release + assets", "~1 min"],
    ["Deploy Demo", "push to main", "build demo + deploy pages", "~2 min"],
]

add_table(slide, Inches(0.5), Inches(4.0), Inches(9.0), Inches(2.0),
          len(workflows), 4, workflows,
          col_widths=[Inches(2.0), Inches(2.5), Inches(3.0), Inches(1.5)])

# Bottom quality gates
add_shape(slide, Inches(0.5), Inches(6.2), Inches(9.0), Inches(0.6), fill_color=DARK_PANEL,
          line_color=TEAL, line_width=Pt(1))
add_textbox(slide, Inches(0.7), Inches(6.25), Inches(8.5), Inches(0.5),
            "🔒 质量门禁: 测试全部通过 + 覆盖率 > 90% + TypeScript零错误 + Lint零警告 → 才允许合并",
            font_size=11, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

print("✓ Slide 27: CI/CD流水线")

# SLIDE 28: 性能指标
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 28)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "性能指标", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Performance dashboard
perf_categories = [
    ("🎮 微游戏", [
        ("单个游戏体积", "< 100KB", TEAL),
        ("游戏生成时间", "< 10ms", GOLD),
        ("帧率", "≥ 60 FPS", ACCENT_RED),
        ("首帧渲染", "< 50ms", TEAL),
    ]),
    ("🚀 启动性能", [
        ("冷启动时间", "< 500ms", ACCENT_RED),
        ("热启动时间", "< 100ms", GOLD),
        ("引擎初始化", "< 200ms", TEAL),
        ("卡片刷新", "< 50ms", RGBColor(0xAA, 0x55, 0xFF)),
    ]),
    ("🤖 AI性能", [
        ("推荐计算", "< 50ms", GOLD),
        ("画像更新", "< 10ms", TEAL),
        ("难度调节", "< 5ms", ACCENT_RED),
        ("触发评估", "< 20ms", RGBColor(0x66, 0xCC, 0x66)),
    ]),
]

for i, (cat_title, metrics) in enumerate(perf_categories):
    x = Inches(0.4) + Inches(3.2) * i
    
    add_shape(slide, x, Inches(1.1), Inches(3.0), Inches(3.3),
              fill_color=DARK_PANEL, line_color=TEAL, line_width=Pt(1))
    add_textbox(slide, x + Inches(0.15), Inches(1.2), Inches(2.7), Inches(0.35),
                cat_title, font_size=14, color=GOLD, bold=True)
    
    for j, (metric, value, color) in enumerate(metrics):
        y = Inches(1.7) + Inches(0.65) * j
        add_textbox(slide, x + Inches(0.15), y, Inches(2.7), Inches(0.25),
                    metric, font_size=10, color=LIGHT_GRAY)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.25), Inches(2.7), Inches(0.3),
                    value, font_size=16, color=color, bold=True)

# Performance comparison
add_textbox(slide, Inches(0.6), Inches(4.7), Inches(8.8), Inches(0.35),
            "📊 对比传统小游戏方案", font_size=16, color=GOLD, bold=True)

comparison_data = [
    ["指标", "传统WebView方案", "Genesis引擎方案", "提升"],
    ["启动时间", "2-5秒", "< 500ms", "4-10x"],
    ["游戏体积", "500KB-5MB", "< 100KB", "5-50x"],
    ["内存占用", "50-200MB", "< 20MB", "2-10x"],
    ["推荐延迟", "500-2000ms (云端)", "< 50ms (本地)", "10-40x"],
    ["离线能力", "受限", "完整支持", "∞"],
]

add_table(slide, Inches(0.5), Inches(5.1), Inches(9.0), Inches(2.1),
          len(comparison_data), 4, comparison_data,
          col_widths=[Inches(2.0), Inches(2.5), Inches(2.5), Inches(1.5)])

print("✓ Slide 28: 性能指标")

# ────────────────────────────────────────────────────────────
# PART 6: 发布计划
# ────────────────────────────────────────────────────────────

# SLIDE 29: Section Divider - Part 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "06", "发布计划", "Release Plan")
add_slide_number(slide, 29)
print("✓ Slide 29: Part 6 Section")

# SLIDE 30: 里程碑路线图
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 30)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "里程碑路线图", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Timeline
phases = [
    ("Phase 0", "基础引擎", "当前 (MVP)", ACCENT_RED, [
        "✅ TypeScript核心引擎",
        "✅ 15个游戏模板",
        "✅ 9大智能触发器",
        "✅ AI画像+推荐系统",
        "✅ 浏览器Demo",
        "✅ 363测试 94%覆盖率",
    ]),
    ("Phase 1", "AI+模板扩展", "Q3 2026", TEAL, [
        "🔲 5+新游戏模板",
        "🔲 深度学习推荐模型",
        "🔲 情绪识别集成",
        "🔲 HarmonyOS Alpha版",
        "🔲 服务卡片实现",
        "🔲 内测用户反馈",
    ]),
    ("Phase 2", "社交+排名", "Q4 2026", GOLD, [
        "🔲 多人对战模式",
        "🔲 全球排行榜",
        "🔲 好友系统",
        "🔲 成就/徽章系统",
        "🔲 HarmonyOS Beta版",
        "🔲 性能深度优化",
    ]),
    ("Phase 3", "开放平台", "Q1 2027", RGBColor(0xAA, 0x55, 0xFF), [
        "🔲 开放API/SDK",
        "🔲 第三方模板市场",
        "🔲 创作者激励计划",
        "🔲 跨平台支持",
        "🔲 HarmonyOS正式版",
        "🔲 商业化运营",
    ]),
]

for i, (phase, title, timeline, color, items) in enumerate(phases):
    x = Inches(0.3) + Inches(2.4) * i
    
    # Phase card
    add_shape(slide, x, Inches(1.1), Inches(2.25), Inches(5.5),
              fill_color=DARK_PANEL, line_color=color, line_width=Pt(2))
    
    # Phase header
    add_shape(slide, x, Inches(1.1), Inches(2.25), Inches(0.9), fill_color=color, shape_type=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, x, Inches(1.12), Inches(2.25), Inches(0.35),
                phase, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(1.45), Inches(2.25), Inches(0.25),
                title, font_size=11, color=RGBColor(0xFF, 0xFF, 0xDD), alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(1.7), Inches(2.25), Inches(0.2),
                timeline, font_size=9, color=RGBColor(0xFF, 0xFF, 0xDD), alignment=PP_ALIGN.CENTER)
    
    # Items
    for j, item in enumerate(items):
        y = Inches(2.2) + Inches(0.4) * j
        item_color = LIGHT_GRAY if item.startswith("🔲") else RGBColor(0x66, 0xCC, 0x66)
        add_textbox(slide, x + Inches(0.1), y, Inches(2.05), Inches(0.35),
                    item, font_size=9, color=item_color)

print("✓ Slide 30: 里程碑路线图")

# SLIDE 31: HarmonyOS上架流程
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 31)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "HarmonyOS 上架流程", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Process steps
upload_steps = [
    ("1", "开发环境", "DevEco Studio\nHarmonyOS SDK\nArkTS开发", ACCENT_RED),
    ("2", "代码迁移", "TS引擎桥接\nArkTS UI实现\n服务卡片开发", TEAL),
    ("3", "AGC配置", "AppGallery Connect\n应用信息注册\n签名配置", GOLD),
    ("4", "测试验证", "真机测试\n兼容性测试\n性能测试", RGBColor(0xAA, 0x55, 0xFF)),
    ("5", "提交审核", "提交应用资料\n隐私政策\n内容审核", ACCENT_RED),
    ("6", "上架发布", "审核通过\n游戏中心上架\n版本管理", RGBColor(0x66, 0xCC, 0x66)),
]

for i, (num, title, desc, color) in enumerate(upload_steps):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(3.15) * col
    y = Inches(1.1) + Inches(2.8) * row
    
    add_shape(slide, x, y, Inches(2.95), Inches(2.5),
              fill_color=DARK_PANEL, line_color=color, line_width=Pt(1.5))
    
    # Number circle
    add_shape(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.5), Inches(0.5),
              fill_color=color, shape_type=MSO_SHAPE.OVAL)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.2), Inches(0.5), Inches(0.4),
                num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_textbox(slide, x + Inches(0.75), y + Inches(0.2), Inches(2.0), Inches(0.35),
                title, font_size=14, color=color, bold=True)
    
    add_textbox(slide, x + Inches(0.15), y + Inches(0.8), Inches(2.6), Inches(1.5),
                desc, font_size=10, color=LIGHT_GRAY)
    
    # Arrow to next
    if i < 5:
        if col < 2:
            add_textbox(slide, x + Inches(2.8), y + Inches(0.8), Inches(0.5), Inches(0.3),
                        "→", font_size=16, color=MEDIUM_GRAY)
        else:
            add_textbox(slide, x + Inches(1.2), y + Inches(2.4), Inches(0.5), Inches(0.3),
                        "↓", font_size=16, color=MEDIUM_GRAY)

print("✓ Slide 31: HarmonyOS上架流程")

# SLIDE 32: 监控与运维
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 32)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "监控与运维", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Four pillars
monitor_pillars = [
    ("📊", "性能监控", ACCENT_RED, [
        "冷启动耗时 < 500ms",
        "游戏帧率 ≥ 60 FPS",
        "AI推荐延迟 < 50ms",
        "内存占用 < 20MB",
        "实时告警阈值",
    ]),
    ("🐛", "崩溃收集", TEAL, [
        "崩溃率 < 0.1%",
        "实时崩溃上报",
        "堆栈追踪还原",
        "用户操作路径回放",
        "自动分级处理",
    ]),
    ("📈", "用户分析", GOLD, [
        "DAU/MAU 统计",
        "游戏偏好分布",
        "触发器命中率",
        "留存率追踪",
        "A/B测试平台",
    ]),
    ("🔄", "热更新", RGBColor(0xAA, 0x55, 0xFF), [
        "游戏模板热更新",
        "AI模型参数更新",
        "触发器规则更新",
        "灰度发布支持",
        "回滚机制保障",
    ]),
]

for i, (icon, title, color, items) in enumerate(monitor_pillars):
    x = Inches(0.4) + Inches(2.4) * i
    
    add_shape(slide, x, Inches(1.1), Inches(2.25), Inches(4.0),
              fill_color=DARK_PANEL, line_color=color, line_width=Pt(1.5))
    
    add_textbox(slide, x + Inches(0.15), Inches(1.2), Inches(0.4), Inches(0.35),
                icon, font_size=20, color=color)
    add_textbox(slide, x + Inches(0.55), Inches(1.22), Inches(1.5), Inches(0.3),
                title, font_size=14, color=color, bold=True)
    
    for j, item in enumerate(items):
        y = Inches(1.8) + Inches(0.42) * j
        add_textbox(slide, x + Inches(0.15), y, Inches(1.95), Inches(0.35),
                    f"▸ {item}", font_size=10, color=LIGHT_GRAY)

# SLA
add_shape(slide, Inches(0.5), Inches(5.4), Inches(9.0), Inches(1.3), fill_color=DARK_PANEL,
          line_color=ACCENT_RED, line_width=Pt(1))
add_textbox(slide, Inches(0.7), Inches(5.5), Inches(8.5), Inches(0.3),
            "📋 SLA 服务等级协议", font_size=13, color=ACCENT_RED, bold=True)

sla_items = [
    ("可用性", "99.9%"),
    ("响应时间", "< 200ms (P99)"),
    ("数据持久性", "99.99%"),
    ("故障恢复", "< 5min"),
]
for i, (label, value) in enumerate(sla_items):
    x = Inches(0.8) + Inches(2.3) * i
    add_textbox(slide, x, Inches(5.85), Inches(1.0), Inches(0.25),
                label, font_size=10, color=LIGHT_GRAY)
    add_textbox(slide, x, Inches(6.1), Inches(1.0), Inches(0.3),
                value, font_size=14, color=GOLD, bold=True)

print("✓ Slide 32: 监控与运维")

# SLIDE 33: 商业化策略
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 33)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "商业化策略", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Revenue model
revenue_models = [
    ("🆓", "免费基础版", "核心变现引擎", ACCENT_RED, [
        "15个基础游戏免费",
        "基础AI推荐功能",
        "每日无限次数",
        "服务卡片完整功能",
        "目标：用户增长 & 留存",
    ]),
    ("📺", "广告收入", "主要收入来源", GOLD, [
        "游戏间隙插屏广告",
        "服务卡片原生广告",
        "激励视频 (复活/道具)",
        "CPM/CPC混合计费",
        "预计: ¥3-5 RPM/用户",
    ]),
    ("👑", "会员订阅", "增值服务", TEAL, [
        "去广告体验",
        "独家游戏模板",
        "高级AI个性化",
        "数据报告与洞察",
        "预计: ¥12/月",
    ]),
    ("💎", "虚拟道具", "辅助变现", RGBColor(0xAA, 0x55, 0xFF), [
        "主题皮肤",
        "特效装饰",
        "排行榜徽章",
        "好友互动道具",
        "预计: ¥1-30/项",
    ]),
]

for i, (icon, title, subtitle, color, items) in enumerate(revenue_models):
    x = Inches(0.4) + Inches(2.4) * i
    
    add_shape(slide, x, Inches(1.1), Inches(2.25), Inches(4.2),
              fill_color=DARK_PANEL, line_color=color, line_width=Pt(1.5))
    
    # Header
    add_textbox(slide, x + Inches(0.1), Inches(1.2), Inches(0.4), Inches(0.35),
                icon, font_size=20, color=color)
    add_textbox(slide, x + Inches(0.5), Inches(1.2), Inches(1.6), Inches(0.3),
                title, font_size=13, color=color, bold=True)
    add_textbox(slide, x + Inches(0.15), Inches(1.55), Inches(1.95), Inches(0.25),
                subtitle, font_size=9, color=MEDIUM_GRAY)
    
    add_accent_line(slide, x + Inches(0.1), Inches(1.85), Inches(2.05), color=color, height=Pt(1))
    
    for j, item in enumerate(items):
        y = Inches(2.0) + Inches(0.4) * j
        add_textbox(slide, x + Inches(0.1), y, Inches(2.05), Inches(0.35),
                    f"▸ {item}", font_size=9, color=LIGHT_GRAY)

# Revenue projection
add_textbox(slide, Inches(0.6), Inches(5.6), Inches(8.8), Inches(0.35),
            "📈 收入预测 (Year 1)", font_size=16, color=GOLD, bold=True)

rev_data = [
    ["指标", "Q1", "Q2", "Q3", "Q4"],
    ["MAU", "50万", "200万", "500万", "1000万"],
    ["付费率", "2%", "3%", "4%", "5%"],
    ["月收入", "¥30万", "¥180万", "¥600万", "¥1500万"],
    ["累计收入", "¥90万", "¥630万", "¥2430万", "¥6930万"],
]

add_table(slide, Inches(0.5), Inches(6.0), Inches(9.0), Inches(1.3),
          len(rev_data), 5, rev_data,
          col_widths=[Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8)])

print("✓ Slide 33: 商业化策略")

# ────────────────────────────────────────────────────────────
# CLOSING
# ────────────────────────────────────────────────────────────

# SLIDE 34: 总结与展望
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_footer(slide, 34)

add_textbox(slide, Inches(0.6), Inches(0.3), Inches(8), Inches(0.5),
            "总结与展望", font_size=26, color=WHITE, bold=True)
add_accent_line(slide, Inches(0.6), Inches(0.85), Inches(1.2), color=ACCENT_RED)

# Key achievements
add_textbox(slide, Inches(0.6), Inches(1.1), Inches(8.8), Inches(0.35),
            "🏆 核心成果", font_size=16, color=GOLD, bold=True)

achievements = [
    ("完整引擎", "15模板 + 9触发器 + AI画像 + 动态难度", "✅"),
    ("高质量交付", "363测试 + 94%覆盖率 + TypeScript严格模式", "✅"),
    ("AI能力闭环", "画像→推荐→反馈→进化 全链路实现", "✅"),
    ("鸿蒙适配方案", "元服务架构 + ArkTS桥接 + 服务卡片", "✅"),
    ("浏览器Demo", "10个可玩游戏 + AI Lab调试面板", "✅"),
]

for i, (title, desc, status) in enumerate(achievements):
    y = Inches(1.5) + Inches(0.55) * i
    add_shape(slide, Inches(0.6), y, Inches(8.8), Inches(0.45),
              fill_color=DARK_PANEL, line_color=RGBColor(0x33, 0x33, 0x55), line_width=Pt(0.5))
    add_shape(slide, Inches(0.6), y, Inches(0.08), Inches(0.45),
              fill_color=RGBColor(0x66, 0xCC, 0x66), shape_type=MSO_SHAPE.RECTANGLE)
    add_textbox(slide, Inches(0.8), y + Inches(0.05), Inches(0.4), Inches(0.3),
                status, font_size=14, color=RGBColor(0x66, 0xCC, 0x66), bold=True)
    add_textbox(slide, Inches(1.3), y + Inches(0.05), Inches(1.5), Inches(0.3),
                title, font_size=12, color=WHITE, bold=True)
    add_textbox(slide, Inches(3.0), y + Inches(0.05), Inches(6.2), Inches(0.3),
                desc, font_size=11, color=LIGHT_GRAY)

# Vision
add_textbox(slide, Inches(0.6), Inches(4.4), Inches(8.8), Inches(0.35),
            "🌟 未来愿景", font_size=16, color=GOLD, bold=True)

visions = [
    ("让游戏无处不在", "每个等待、每个碎片时间，都有专属游戏陪伴"),
    ("AI驱动的游戏生态", "千人千面的游戏体验，越玩越懂你"),
    ("鸿蒙原生游戏引擎", "成为HarmonyOS小游戏的事实标准"),
    ("开放平台战略", "赋能开发者，共建微游戏生态"),
]

for i, (title, desc) in enumerate(visions):
    y = Inches(4.9) + Inches(0.55) * i
    add_textbox(slide, Inches(0.8), y, Inches(2.5), Inches(0.3),
                f"▸ {title}", font_size=12, color=TEAL, bold=True)
    add_textbox(slide, Inches(3.5), y, Inches(6), Inches(0.3),
                desc, font_size=11, color=LIGHT_GRAY)

print("✓ Slide 34: 总结与展望")

# SLIDE 35: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Top accent bar
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), fill_color=ACCENT_RED, shape_type=MSO_SHAPE.RECTANGLE)

# Genesis logo
add_shape(slide, Inches(3.5), Inches(1.2), Inches(3), Inches(0.6), fill_color=None,
          line_color=ACCENT_RED, line_width=Pt(2), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, Inches(3.5), Inches(1.25), Inches(3), Inches(0.5),
            "◆  GENESIS  ◆", font_size=20, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)

# Thank you
add_textbox(slide, Inches(1), Inches(2.3), Inches(8), Inches(1.0),
            "Thank You", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Chinese subtitle
add_textbox(slide, Inches(1), Inches(3.3), Inches(8), Inches(0.6),
            "感谢聆听", font_size=28, color=GOLD, alignment=PP_ALIGN.CENTER)

# Accent line
add_accent_line(slide, Inches(3.5), Inches(4.1), Inches(3), color=ACCENT_RED, height=Pt(2))

# Tagline
add_textbox(slide, Inches(1), Inches(4.5), Inches(8), Inches(0.5),
            "让每个等待时刻都变成游戏时刻", font_size=18, color=TEAL, alignment=PP_ALIGN.CENTER)

# Contact info
contact_items = [
    "📧  genesis-engine@example.com",
    "🌐  github.com/genesis-ai-micro-game-engine",
    "📱  华为游戏中心 × 鸿蒙元服务",
]
for i, contact in enumerate(contact_items):
    y = Inches(5.3) + Inches(0.4) * i
    add_textbox(slide, Inches(1), y, Inches(8), Inches(0.35),
                contact, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom version
add_textbox(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.4),
            "Genesis AI Micro-Game Engine  |  Version 2.0 MVP  |  © 2026",
            font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

print("✓ Slide 35: Thank You")

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = "/root/genesis/docs/Genesis_AI_微游戏生成器_方案升级版.pptx"
prs.save(output_path)
print(f"\n{'='*60}")
print(f"✅ Presentation saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"{'='*60}")
